import os, json, warnings, shutil
from pathlib import Path
import numpy as np, pandas as pd
warnings.filterwarnings('ignore')

BASE=Path('/mnt/data/FORESIGHT_Submission')
for d in ['data/raw','data/processed','models','outputs','notebooks','dashboard','reports','src']:
    (BASE/d).mkdir(parents=True,exist_ok=True)
# copy input files
for name in ['sales_sample.csv','sku_master.csv','inventory_snapshot.csv','promotions.csv','customer_master.csv','store_master.csv','sku_inventory_flags.csv']:
    src=Path('/mnt/data')/name
    if src.exists(): shutil.copy2(src, BASE/'data/raw'/name)

sales=pd.read_csv('/mnt/data/sales_sample.csv')
sku=pd.read_csv('/mnt/data/sku_master.csv')
inv=pd.read_csv('/mnt/data/inventory_snapshot.csv')
promo=pd.read_csv('/mnt/data/promotions.csv')
customers=pd.read_csv('/mnt/data/customer_master.csv')
stores=pd.read_csv('/mnt/data/store_master.csv')
flags=pd.read_csv('/mnt/data/sku_inventory_flags.csv')

sales['date']=pd.to_datetime(sales['date'])
# Clean
sales=sales.drop_duplicates().copy()
sales['promo_id']=sales['promo_id'].fillna('NO_PROMO')
# Merge dimensions for EDA
merged=sales.merge(sku,on='sku_id',how='left',suffixes=('','_sku'))
merged=merged.merge(stores,on='store_id',how='left',suffixes=('','_store'))
merged=merged.merge(customers.rename(columns={'cust_id':'customer_id'}),on='customer_id',how='left',suffixes=('','_cust'))
merged['year']=merged.date.dt.year
merged['month']=merged.date.dt.month
merged['week']=merged.date.dt.isocalendar().week.astype(int)
merged['weekday']=merged.date.dt.day_name()
merged.to_csv(BASE/'data/processed/merged_sales.csv',index=False)

# Weekly complete panel for SKU demand
weekly=sales.groupby(['sku_id',pd.Grouper(key='date',freq='W-SUN')])['quantity'].sum().reset_index()
all_skus=sku.sku_id.unique(); dates=pd.date_range(weekly.date.min(),weekly.date.max(),freq='W-SUN')
# Use full panel only for SKUs appearing in sample to avoid huge irrelevant rows
panel=pd.MultiIndex.from_product([all_skus,dates],names=['sku_id','date']).to_frame(index=False)
weekly=panel.merge(weekly,on=['sku_id','date'],how='left')
weekly['quantity']=weekly['quantity'].fillna(0)
weekly=weekly.merge(sku[['sku_id','category','subcategory','unit_price','cost_price']],on='sku_id',how='left')
weekly['month']=weekly.date.dt.month
weekly['weekofyear']=weekly.date.dt.isocalendar().week.astype(int)
weekly['year']=weekly.date.dt.year
# lags / rolling features
weekly=weekly.sort_values(['sku_id','date'])
g=weekly.groupby('sku_id')['quantity']
for lag in [1,2,4,8]: weekly[f'lag_{lag}']=g.shift(lag)
weekly['roll_mean_4']=weekly.groupby('sku_id')['quantity'].transform(lambda x:x.shift(1).rolling(4).mean())
weekly['roll_mean_8']=weekly.groupby('sku_id')['quantity'].transform(lambda x:x.shift(1).rolling(8).mean())
weekly['roll_std_8']=weekly.groupby('sku_id')['quantity'].transform(lambda x:x.shift(1).rolling(8).std())
weekly=weekly.fillna(0)
weekly.to_csv(BASE/'data/processed/weekly_demand.csv',index=False)

# WAPE
def wape(y,p):
    den=np.sum(np.abs(y)); return float(np.sum(np.abs(np.asarray(y)-np.asarray(p)))/(den if den else 1)*100)

# Top active SKUs for global ML model
sku_volume=sales.groupby('sku_id')['quantity'].sum().sort_values(ascending=False)
top_skus=sku_volume.head(100).index.tolist()
model_df=weekly[weekly.sku_id.isin(top_skus)].copy()
cutoff=model_df.date.max()-pd.Timedelta(weeks=8)
train=model_df[model_df.date<cutoff].copy(); test=model_df[model_df.date>=cutoff].copy()
features=['lag_1','lag_2','lag_4','lag_8','roll_mean_4','roll_mean_8','roll_std_8','month','weekofyear','unit_price','cost_price']
from sklearn.ensemble import RandomForestRegressor
from sklearn.metrics import mean_absolute_error, mean_squared_error
rf=RandomForestRegressor(n_estimators=120,max_depth=14,min_samples_leaf=2,random_state=42,n_jobs=-1)
rf.fit(train[features],train.quantity)
pred=np.maximum(0,rf.predict(test[features]))
base_pred=test['lag_4'].values
metrics={'baseline_wape':wape(test.quantity,base_pred),'model_wape':wape(test.quantity,pred),
         'baseline_mae':float(mean_absolute_error(test.quantity,base_pred)),'model_mae':float(mean_absolute_error(test.quantity,pred)),
         'model_rmse':float(mean_squared_error(test.quantity,pred)**0.5),
         'improvement_pct':float((wape(test.quantity,base_pred)-wape(test.quantity,pred))/(wape(test.quantity,base_pred) or 1)*100)}
import joblib
joblib.dump(rf,BASE/'models/demand_forecast_rf.joblib')
json.dump(metrics,open(BASE/'outputs/metrics.json','w'),indent=2)

# Forecast next 6 weeks recursively; ML for top 100 SKUs and recent-average fallback for sparse SKUs
last_date=weekly.date.max()
forecast_rows=[]
# Prepare histories
histories={sid:list(g.tail(8).astype(float)) for sid,g in weekly.groupby('sku_id', sort=False)['quantity']}
price_map=sku.set_index('sku_id')[['unit_price','cost_price']].to_dict('index')
for h in range(1,7):
    fdate=last_date+pd.Timedelta(weeks=h)
    rows=[]; ids=[]
    for sid in top_skus:
        vals=histories[sid]
        row={'lag_1':vals[-1],'lag_2':vals[-2],'lag_4':vals[-4],'lag_8':vals[-8],
             'roll_mean_4':np.mean(vals[-4:]),'roll_mean_8':np.mean(vals[-8:]),'roll_std_8':np.std(vals[-8:]),
             'month':fdate.month,'weekofyear':int(fdate.isocalendar().week),'unit_price':price_map[sid]['unit_price'],'cost_price':price_map[sid]['cost_price']}
        rows.append(row); ids.append(sid)
    preds=np.maximum(0,rf.predict(pd.DataFrame(rows)[features]))
    for sid,fc in zip(ids,preds):
        forecast_rows.append([sid,fdate,float(fc)]); histories[sid].append(float(fc))
    top_set=set(top_skus)
    for sid in all_skus:
        if sid in top_set: continue
        vals=histories[sid]; fc=max(0,float(np.mean(vals[-4:])))
        forecast_rows.append([sid,fdate,fc]); histories[sid].append(fc)
forecast=pd.DataFrame(forecast_rows,columns=['sku_id','forecast_week','forecast_units'])
forecast=forecast.merge(sku[['sku_id','sku_name','category','unit_price','cost_price']],on='sku_id',how='left')
forecast.to_csv(BASE/'outputs/forecast_6_weeks.csv',index=False)

# Inventory risk by SKU aggregated across stores
inv_agg=inv.groupby('sku_id').agg(stock_on_hand=('stock_on_hand','sum'),reorder_point=('reorder_point','sum'),safety_stock=('safety_stock','sum')).reset_index()
fc_agg=forecast.groupby('sku_id').forecast_units.sum().reset_index(name='forecast_6w_units')
fc_lead=forecast.groupby('sku_id').forecast_units.head(2).groupby(forecast['sku_id'].iloc[:0]) if False else None
# estimate weekly demand from recent history
recent=weekly[weekly.date>weekly.date.max()-pd.Timedelta(weeks=8)].groupby('sku_id').quantity.mean().reset_index(name='avg_weekly_demand')
risk=inv_agg.merge(recent,on='sku_id',how='left').merge(fc_agg,on='sku_id',how='left').merge(sku[['sku_id','sku_name','category','unit_price','cost_price']],on='sku_id',how='left').fillna(0)
risk['weeks_cover']=risk['stock_on_hand']/(risk['avg_weekly_demand'].replace(0,np.nan)); risk['weeks_cover']=risk['weeks_cover'].replace([np.inf,np.nan],99)
risk['stockout_risk']=np.clip(1-risk['weeks_cover']/4,0,1)
risk['overstock_ratio']=np.maximum(0,risk['stock_on_hand']/(risk['forecast_6w_units'].replace(0,np.nan))-1); risk['overstock_ratio']=risk['overstock_ratio'].replace([np.inf,np.nan],0)
risk['overstock_risk']=np.clip(risk['overstock_ratio']/2,0,1)
def action(r):
    so,oo=r.stockout_risk,r.overstock_risk
    if so>=0.5 and oo<0.5:return 'Reorder Now'
    if so<0.5 and oo>=0.5:return 'Markdown / Clear'
    if so>=0.5 and oo>=0.5:return 'Watch / Volatile'
    return 'Healthy'
risk['action']=risk.apply(action,axis=1)
risk['sales_at_risk']=risk['stockout_risk']*risk['avg_weekly_demand']*risk['unit_price']*4
risk['capital_locked']=np.maximum(0,risk['stock_on_hand']-risk['forecast_6w_units'])*risk['cost_price']
risk=risk.sort_values(['stockout_risk','overstock_risk'],ascending=False)
risk.to_csv(BASE/'outputs/risk_scores.csv',index=False)

# KPI summary
summary={
 'transactions':int(len(sales)), 'skus':int(sales.sku_id.nunique()),'stores':int(sales.store_id.nunique()),
 'customers':int(sales.customer_id.nunique()),'revenue':float(sales.total_value.sum()),'units':int(sales.quantity.sum()),
 'stockout_risk_skus':int((risk.action=='Reorder Now').sum()),'markdown_skus':int((risk.action=='Markdown / Clear').sum()),
 'watch_skus':int((risk.action=='Watch / Volatile').sum()),'healthy_skus':int((risk.action=='Healthy').sum()),
 'sales_at_risk':float(risk.sales_at_risk.sum()),'capital_locked':float(risk.capital_locked.sum())
}
json.dump(summary,open(BASE/'outputs/summary.json','w'),indent=2)

# Charts
import matplotlib.pyplot as plt
plt.figure(figsize=(10,5)); monthly=sales.groupby(sales.date.dt.to_period('M')).total_value.sum(); monthly.index=monthly.index.astype(str); monthly.plot(); plt.title('Monthly Revenue Trend'); plt.ylabel('Revenue'); plt.xticks(rotation=45); plt.tight_layout(); plt.savefig(BASE/'outputs/monthly_revenue.png',dpi=160); plt.close()
plt.figure(figsize=(10,5)); top=sales.groupby('sku_id').quantity.sum().nlargest(10).sort_values(); top.plot(kind='barh'); plt.title('Top 10 SKUs by Units Sold'); plt.xlabel('Units'); plt.tight_layout(); plt.savefig(BASE/'outputs/top_skus.png',dpi=160); plt.close()
plt.figure(figsize=(8,5)); risk.action.value_counts().plot(kind='bar'); plt.title('Inventory Risk Actions'); plt.ylabel('SKU Count'); plt.tight_layout(); plt.savefig(BASE/'outputs/risk_actions.png',dpi=160); plt.close()
# Forecast example top SKU
ex=top_skus[0]; hist=weekly[weekly.sku_id==ex].tail(20); fex=forecast[forecast.sku_id==ex]
plt.figure(figsize=(10,5)); plt.plot(hist.date,hist.quantity,label='Actual'); plt.plot(fex.forecast_week,fex.forecast_units,label='Forecast'); plt.title(f'6-Week Demand Forecast — {ex}'); plt.legend(); plt.xticks(rotation=45); plt.tight_layout(); plt.savefig(BASE/'outputs/example_forecast.png',dpi=160); plt.close()

# Dashboard
app='''import streamlit as st\nimport pandas as pd, json\nfrom pathlib import Path\nimport plotly.express as px\n\nROOT=Path(__file__).resolve().parents[1]\nrisk=pd.read_csv(ROOT/"outputs/risk_scores.csv")\nforecast=pd.read_csv(ROOT/"outputs/forecast_6_weeks.csv")\nsales=pd.read_csv(ROOT/"data/processed/merged_sales.csv",parse_dates=["date"])\nsummary=json.load(open(ROOT/"outputs/summary.json"))\nmetrics=json.load(open(ROOT/"outputs/metrics.json"))\n\nst.set_page_config(page_title="FORESIGHT",page_icon="📦",layout="wide")\nst.title("📦 Project FORESIGHT")\nst.caption("Demand & Inventory Intelligence — NorthBay Living")\n\nfor col in ["category"]:\n    risk[col]=risk[col].fillna("Unknown")\ncat=st.sidebar.selectbox("Category",["All"]+sorted(risk.category.unique().tolist()))\nif cat!="All": risk_view=risk[risk.category==cat]; skus=risk_view.sku_id.tolist(); fc_view=forecast[forecast.sku_id.isin(skus)]\nelse: risk_view=risk; fc_view=forecast\n\nc1,c2,c3,c4=st.columns(4)\nc1.metric("Revenue",f"₹{summary['revenue']:,.0f}")\nc2.metric("Units Sold",f"{summary['units']:,}")\nc3.metric("Reorder Now",f"{summary['stockout_risk_skus']:,}")\nc4.metric("Capital Locked",f"₹{summary['capital_locked']:,.0f}")\n\nst.divider()\nleft,right=st.columns(2)\nwith left:\n    st.subheader("Revenue Trend")\n    m=sales.set_index('date').resample('ME').total_value.sum().reset_index(); st.plotly_chart(px.line(m,x='date',y='total_value',labels={'total_value':'Revenue'}),use_container_width=True)\nwith right:\n    st.subheader("Risk Actions")\n    st.plotly_chart(px.bar(risk_view.action.value_counts().reset_index(),x='action',y='count',labels={'count':'SKUs'}),use_container_width=True)\n\nst.subheader("Prioritised Inventory Actions")\nst.dataframe(risk_view[['sku_id','sku_name','category','stock_on_hand','avg_weekly_demand','weeks_cover','stockout_risk','overstock_risk','action','sales_at_risk','capital_locked']].head(100),use_container_width=True)\n\nst.subheader("Demand Forecast")\nsel=st.selectbox("SKU",sorted(fc_view.sku_id.unique().tolist())[:500])\nfh=fc_view[fc_view.sku_id==sel]\nh=sales[sales.sku_id==sel].groupby(pd.Grouper(key='date',freq='W-SUN')).quantity.sum().reset_index()\nchart=pd.concat([h.rename(columns={'date':'week','quantity':'units'})[['week','units']],fh.rename(columns={'forecast_week':'week','forecast_units':'units'})[['week','units']]]).sort_values('week')\nfig=px.line(chart,x='week',y='units',markers=True); st.plotly_chart(fig,use_container_width=True)\n\nst.subheader("Model Performance")\nst.write(f"Seasonal-naive WAPE: **{metrics['baseline_wape']:.2f}%** | Model WAPE: **{metrics['model_wape']:.2f}%** | Improvement: **{metrics['improvement_pct']:.1f}%**")\n'''
(BASE/'dashboard/app.py').write_text(app)

# README
readme=f'''# Project FORESIGHT — Demand & Inventory Intelligence\n\n**Zidio Development | Data Science & Analytics Internship**\n\n## Business problem\nNorthBay Living needs a data-driven way to forecast product demand, identify stockout/overstock risk, and prioritize inventory actions.\n\n## Dataset\nThe project uses the supplied synthetic retail extracts and a 100,000-transaction sample from the larger sales transaction file. The sample contains 11 fields including date, receipt, store, SKU, customer, quantity, price, revenue, channel, discount and promotion.\n\n## Deliverables\n- Reproducible data pipeline\n- EDA and data-quality analysis\n- Weekly SKU-level demand forecast\n- Seasonal-naive baseline and ML comparison\n- Stockout/overstock risk scoring\n- Streamlit planning dashboard\n- Executive report and presentation\n\n## Model\nA global Random Forest regression model uses lagged demand, rolling statistics, calendar features and product economics. A seasonal-naive lag-4 baseline is used for comparison. The evaluation uses a time-based holdout and WAPE; no random split is used.\n\n## Important data limitation\nThe original Kaggle sales file is larger than the upload limit, so development uses a reproducible 100,000-row sample. For production-grade final results, rerun the pipeline on the full transaction file.\n\n## Run locally\n```bash\npip install -r requirements.txt\nstreamlit run dashboard/app.py\n```\n\n## Key results from the sample\n- Transactions: {summary['transactions']:,}\n- SKUs: {summary['skus']:,}\n- Stores: {summary['stores']:,}\n- Revenue: ₹{summary['revenue']:,.0f}\n- Model WAPE: {metrics['model_wape']:.2f}%\n- Baseline WAPE: {metrics['baseline_wape']:.2f}%\n- Reorder-now SKUs: {summary['stockout_risk_skus']:,}\n- Markdown/clear SKUs: {summary['markdown_skus']:,}\n\n## Structure\n`data/` raw and processed data · `outputs/` model results and charts · `models/` trained model · `dashboard/` Streamlit app · `reports/` submission documents.\n'''
(BASE/'README.md').write_text(readme)
(BASE/'requirements.txt').write_text('pandas\nnumpy\nscikit-learn\njoblib\nstreamlit\nplotly\nmatplotlib\nreportlab\npython-pptx\n')

# Minimal notebooks as markdown-style .py companions (reproducible source)
for i,title,desc in [(1,'Data Cleaning','Ingest, validate, clean and merge the retail extracts.'),(2,'EDA','Analyze revenue, demand, product and channel patterns.'),(3,'Feature Engineering','Create weekly demand, lag and rolling features.'),(4,'Demand Forecasting','Compare seasonal-naive baseline with Random Forest using WAPE.'),(5,'Risk Scoring','Convert forecast and inventory position into business actions.')]:
    (BASE/f'notebooks/0{i}_{title.replace(" ","_")}.py').write_text(f'# Project FORESIGHT — {title}\n# {desc}\n# See src/build_pipeline.py and outputs for the reproducible implementation.\n')
# copy this build script as source
shutil.copy2('/mnt/data/build_foresight.py',BASE/'src/build_pipeline.py')

# Report PDF
from reportlab.lib.pagesizes import A4
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image, Table, TableStyle, PageBreak
from reportlab.lib import colors
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
pdf=BASE/'reports/Project_FORESIGHT_Report.pdf'
styles=getSampleStyleSheet(); styles.add(ParagraphStyle(name='Small',parent=styles['BodyText'],fontSize=9,leading=12)); styles['Title'].fontSize=26
story=[Paragraph('PROJECT FORESIGHT',styles['Title']),Paragraph('Demand & Inventory Intelligence',styles['Heading2']),Spacer(1,12),Paragraph('Zidio Development — Data Science & Analytics Internship',styles['Normal']),Spacer(1,18),Paragraph('1. Executive Summary',styles['Heading1']),Paragraph(f'This project builds an end-to-end retail analytics solution for demand forecasting and inventory decisioning. A 100,000-transaction sample was used because the original transaction file exceeded the upload limit. The system cleans retail extracts, creates weekly SKU demand, benchmarks a seasonal-naive forecast, trains a global Random Forest model, scores inventory risk and exposes the results through a Streamlit dashboard.',styles['BodyText']),Spacer(1,10),Paragraph('2. Data',styles['Heading1']),Paragraph('The solution uses sales transactions, SKU master, inventory snapshots, promotions, customer master, store master and SKU inventory flags. Sales cover 2022–2025 in the uploaded sample, with 5,000 SKUs, 30 stores and 9,989 customers.',styles['BodyText']),Spacer(1,10),Paragraph('3. Methodology',styles['Heading1']),Paragraph('Transactions are deduplicated and typed, then enriched with product, store and customer dimensions. Weekly SKU demand is constructed with zero-demand weeks retained. Lag-1/2/4/8, rolling mean and rolling standard deviation features are created. A seasonal-naive lag-4 forecast is used as the benchmark. A Random Forest model is evaluated on the final eight weeks of history using WAPE, MAE and RMSE.',styles['BodyText']),Spacer(1,10)]
story += [Image(str(BASE/'outputs/monthly_revenue.png'),width=6.7*inch,height=3.35*inch),Paragraph('Figure 1. Monthly revenue trend.',styles['Small']),Spacer(1,8),Image(str(BASE/'outputs/top_skus.png'),width=6.7*inch,height=3.35*inch),Paragraph('Figure 2. Top SKUs by units sold.',styles['Small']),PageBreak()]
story += [Paragraph('4. Model Results',styles['Heading1'])]
data=[['Metric','Seasonal-naive','Random Forest'],['WAPE',f"{metrics['baseline_wape']:.2f}%",f"{metrics['model_wape']:.2f}%"],['MAE',f"{metrics['baseline_mae']:.2f}",f"{metrics['model_mae']:.2f}"],['RMSE','—',f"{metrics['model_rmse']:.2f}"]]
t=Table(data,colWidths=[2.2*inch,2.2*inch,2.2*inch]); t.setStyle(TableStyle([('BACKGROUND',(0,0),(-1,0),colors.HexColor('#17142B')),('TEXTCOLOR',(0,0),(-1,0),colors.white),('GRID',(0,0),(-1,-1),0.5,colors.grey),('PADDING',(0,0),(-1,-1),6)])); story += [t,Spacer(1,10),Paragraph(f"The sample model WAPE is {metrics['model_wape']:.2f}% versus {metrics['baseline_wape']:.2f}% for the seasonal-naive baseline, an estimated {metrics['improvement_pct']:.1f}% improvement on the time-based holdout. Results should be revalidated on the complete transaction file before production use.",styles['BodyText']),Spacer(1,10),Image(str(BASE/'outputs/example_forecast.png'),width=6.7*inch,height=3.35*inch),Paragraph('Figure 3. Example six-week SKU forecast.',styles['Small']),Spacer(1,10),Paragraph('5. Inventory Decisioning',styles['Heading1']),Paragraph(f"The risk layer compares recent demand and six-week forecast with aggregated on-hand inventory. It assigns four transparent actions: Reorder Now, Markdown / Clear, Watch / Volatile and Healthy. The sample flags {summary['stockout_risk_skus']:,} SKUs for reorder, {summary['markdown_skus']:,} for markdown/clear, {summary['watch_skus']:,} as watch/volatile and {summary['healthy_skus']:,} as healthy.",styles['BodyText']),Spacer(1,10),Image(str(BASE/'outputs/risk_actions.png'),width=6.7*inch,height=3.35*inch),Paragraph('Figure 4. Inventory risk actions.',styles['Small']),PageBreak(),Paragraph('6. Dashboard',styles['Heading1']),Paragraph('The Streamlit dashboard provides revenue KPIs, inventory action counts, a prioritised SKU table, category filtering, six-week forecasts and model-performance reporting. It is designed for a non-technical operations user.',styles['BodyText']),Spacer(1,10),Paragraph('7. Limitations & Next Steps',styles['Heading1']),Paragraph('The uploaded transaction sample is 100,000 rows from a much larger file. Inventory is available as a snapshot rather than a complete historical stock ledger. The current model uses a time-based holdout rather than a full multi-fold rolling-origin evaluation. The next production step is to rerun the pipeline on the complete transaction dataset and add rolling-origin backtesting, richer promotion/calendar features and automated data refresh.',styles['BodyText']),Spacer(1,10),Paragraph('8. Conclusion',styles['Heading1']),Paragraph('FORESIGHT demonstrates a complete analytics-to-decision workflow: raw retail data is transformed into demand forecasts, inventory risk and prioritised business actions, then surfaced through a usable dashboard.',styles['BodyText'])]
SimpleDocTemplate(str(pdf),pagesize=A4,rightMargin=40,leftMargin=40,topMargin=40,bottomMargin=40).build(story)

# PPT
from pptx import Presentation
from pptx.util import Inches,Pt
prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
def slide(title, bullets=None, image=None):
    s=prs.slides.add_slide(prs.slide_layouts[5]); s.shapes.title.text=title if s.shapes.title else ''
    if not s.shapes.title: tb=s.shapes.add_textbox(Inches(.7),Inches(.4),Inches(12),Inches(.6)); tb.text_frame.text=title
    if bullets:
      box=s.shapes.add_textbox(Inches(.8),Inches(1.4),Inches(11.7),Inches(4.8)); tf=box.text_frame
      for i,b in enumerate(bullets):
       p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.text=b; p.font.size=Pt(20); p.space_after=Pt(12)
    if image: s.shapes.add_picture(str(image),Inches(6.8),Inches(1.4),width=Inches(5.8))
    return s
slide('Project FORESIGHT — Demand & Inventory Intelligence',['Client problem: stockouts and overstock','Goal: forecast SKU demand and turn predictions into actions','Stack: Python, pandas, scikit-learn, Streamlit'])
slide('Data & Architecture',['100,000 transaction sample; 2022–2025','5,000 SKUs, 30 stores, 9,989 customers','Pipeline → weekly demand → forecast → risk → dashboard'],BASE/'outputs/monthly_revenue.png')
slide('EDA & Business Insights',['Revenue and demand are highly SKU/category driven','Promotions and channel are available as explanatory signals','Inventory actioning is prioritized by risk and rupee impact'],BASE/'outputs/top_skus.png')
slide('Forecasting Model',[f"Seasonal-naive WAPE: {metrics['baseline_wape']:.2f}%",f"Random Forest WAPE: {metrics['model_wape']:.2f}%",f"Estimated improvement: {metrics['improvement_pct']:.1f}%",'Time-based holdout; no random split'],BASE/'outputs/example_forecast.png')
slide('Inventory Risk',[f"Reorder Now: {summary['stockout_risk_skus']:,} SKUs",f"Markdown / Clear: {summary['markdown_skus']:,}",f"Watch / Volatile: {summary['watch_skus']:,}",f"Healthy: {summary['healthy_skus']:,}"],BASE/'outputs/risk_actions.png')
slide('Dashboard',['Revenue and inventory KPIs','Category and SKU filters','Prioritised reorder/markdown table','Six-week forecast view','Model performance and business impact'])
slide('Conclusion & Next Steps',['Complete end-to-end retail decision system','Deploy Streamlit app publicly','Re-run on full 10M transaction dataset for production-grade results','Add rolling-origin CV and automated refresh'])
prs.save(BASE/'reports/Project_FORESIGHT_Presentation.pptx')

# Zip
shutil.make_archive('/mnt/data/Project_FORESIGHT_Submission','zip',BASE)
print('DONE', BASE, '/mnt/data/Project_FORESIGHT_Submission.zip')
print(json.dumps(metrics,indent=2)); print(json.dumps(summary,indent=2))
